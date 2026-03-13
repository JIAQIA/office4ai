"""
PPT E2E 测试基础设施

提供 PowerPoint 专属的文档读取器、AppleScript 函数和测试运行器。

使用方式:
    from manual_tests.ppt.e2e_base import PPTTestRunner, PresentationReader

    runner = PPTTestRunner()
    async with runner.run_with_workspace("empty.pptx") as (workspace, fixture):
        result = await workspace.execute(action)
        reader = PresentationReader(fixture.working_path)
        reader.reload()
        assert reader.contains_text("Hello PPT")
"""

import asyncio
import inspect
import platform
import subprocess
from collections.abc import AsyncIterator, Callable
from contextlib import asynccontextmanager
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from pptx import Presentation as PptxPresentation

from manual_tests.e2e_base import (
    DocumentFixture,
    E2ETestRunner,
    TestCase,
    open_document,
    path_to_file_uri,
)
from office4ai.environment.workspace.office_workspace import OfficeWorkspace

# ==============================================================================
# 配置常量
# ==============================================================================

DEFAULT_PPT_ADDIN_NAME = "ppt-editor"

FIXTURES_ROOT = Path(__file__).parent / "fixtures"

TEMP_ROOT = Path(__file__).parent / ".test_working"

# ==============================================================================
# Validator 类型定义
# ==============================================================================

# 传统验证器：仅接收协议返回的 data
PptDataValidator = Callable[[dict[str, Any]], bool]

# 内容验证器：接收 data 和 PresentationReader，用于双重验证
PptContentValidator = Callable[[dict[str, Any], "PresentationReader"], bool]

# 统一验证器类型
PptValidator = PptDataValidator | PptContentValidator


def _call_ppt_validator(
    validator: PptValidator,
    data: dict[str, Any],
    reader: "PresentationReader",
) -> bool:
    """
    智能调用验证器

    根据验证器的参数数量自动判断是传统模式还是双重验证模式。
    """
    sig = inspect.signature(validator)
    if len(sig.parameters) == 2:
        return validator(data, reader)  # type: ignore[call-arg]
    else:
        return validator(data)  # type: ignore[call-arg]


# ==============================================================================
# PresentationReader
# ==============================================================================


@dataclass
class PresentationReader:
    """
    PPT 文档内容读取器

    提供对 PowerPoint 文档内容的只读访问，用于双重验证。
    通过 python-pptx 读取修改后的文档，验证实际内容。

    Example:
        reader = PresentationReader(fixture.working_path)
        reader.reload()
        assert reader.contains_text("Hello PPT")
        assert reader.slide_has_text(0, "First Slide")
    """

    path: Path
    auto_save: bool = True
    _prs: "PptxPresentation | None" = field(default=None, repr=False)

    @property
    def prs(self) -> "PptxPresentation":
        """懒加载 Presentation 对象"""
        if self._prs is None:
            from pptx import Presentation

            self._prs = Presentation(str(self.path))
        return self._prs

    def reload(self) -> None:
        """
        重新加载演示文稿

        当 auto_save=True 时，先通过 AppleScript 强制 PowerPoint 保存文档到磁盘，
        再重新加载。
        """
        if self.auto_save:
            save_ppt_document(self.path)
        self._prs = None

    @property
    def slide_count(self) -> int:
        """获取幻灯片数量"""
        return len(self.prs.slides)

    def get_all_text(self) -> str:
        """获取演示文稿中所有文本（所有幻灯片的文本用换行连接）"""
        texts: list[str] = []
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
        return "\n".join(texts)

    def contains_text(self, text: str) -> bool:
        """检查演示文稿中是否包含指定文本"""
        return text in self.get_all_text()

    def slide_has_text(self, index: int, text: str) -> bool:
        """
        检查指定幻灯片是否包含文本

        Args:
            index: 幻灯片索引（0-based）
            text: 要搜索的文本
        """
        slides = list(self.prs.slides)
        if 0 <= index < len(slides):
            slide = slides[index]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if text in paragraph.text:
                            return True
        return False

    def get_slide_shapes(self, index: int) -> list[Any]:
        """
        获取指定幻灯片的所有形状

        Args:
            index: 幻灯片索引（0-based）

        Returns:
            形状列表
        """
        slides = list(self.prs.slides)
        if 0 <= index < len(slides):
            return list(slides[index].shapes)
        return []

    def table_count(self, slide_index: int) -> int:
        """
        获取指定幻灯片上的表格数量

        Args:
            slide_index: 幻灯片索引（0-based）
        """
        count = 0
        for shape in self.get_slide_shapes(slide_index):
            if shape.has_table:
                count += 1
        return count

    def get_table_cell_text(self, slide_index: int, table_index: int, row: int, col: int) -> str | None:
        """
        获取表格单元格文本

        Args:
            slide_index: 幻灯片索引（0-based）
            table_index: 表格索引（0-based，在该幻灯片上）
            row: 行索引（0-based）
            col: 列索引（0-based）

        Returns:
            单元格文本，如果找不到则返回 None
        """
        table_idx = 0
        for shape in self.get_slide_shapes(slide_index):
            if shape.has_table:
                if table_idx == table_index:
                    table = shape.table
                    if 0 <= row < len(table.rows) and 0 <= col < len(table.columns):
                        return table.cell(row, col).text
                    return None
                table_idx += 1
        return None


# ==============================================================================
# 预期数据类
# ==============================================================================


@dataclass
class ExpectedSlideInfo:
    """
    预期的幻灯片信息

    用于自动验证 ppt:get:slideInfo 结果。
    设置为 None 的字段将跳过验证。
    """

    slide_count: int | None = None
    current_index: int | None = None
    width: float | None = None
    height: float | None = None

    # 容差
    slide_count_tolerance: int = 0
    width_tolerance: float = 1.0
    height_tolerance: float = 1.0


@dataclass
class PptTestCase:
    """
    PPT 测试用例定义

    Attributes:
        name: 测试名称
        fixture_name: 夹具文件名（相对于 fixtures 目录中的 PPT 子目录）
        description: 测试描述
        expected: 预期的幻灯片信息（可选）
        validator: 自定义验证函数（可选）
        tags: 标签列表
    """

    name: str
    fixture_name: str
    description: str
    expected: ExpectedSlideInfo | None = None
    validator: PptValidator | None = None
    tags: list[str] = field(default_factory=list)


# ==============================================================================
# AppleScript 函数
# ==============================================================================


def _run_applescript(script: str, timeout: float = 15.0) -> tuple[bool, str]:
    """执行 AppleScript 并返回结果"""
    try:
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True,
            timeout=timeout,
        )
        output = result.stdout.decode().strip()
        if result.returncode == 0:
            return True, output
        else:
            stderr = result.stderr.decode().strip()
            return False, stderr
    except subprocess.TimeoutExpired:
        return False, "AppleScript 执行超时"
    except Exception as e:
        return False, str(e)


def save_ppt_document(path: Path) -> bool:
    """
    强制 PowerPoint 保存文档（通过 AppleScript）

    Args:
        path: 文档路径

    Returns:
        是否成功
    """
    if platform.system() != "Darwin":
        print("⚠️  自动保存仅支持 macOS")
        return False

    try:
        doc_name = path.name
        doc_stem = path.stem
        script = f'''
        tell application "Microsoft PowerPoint"
            set targetDocs to (every presentation whose name is "{doc_name}")
            if (count of targetDocs) is 0 then
                set targetDocs to (every presentation whose name starts with "{doc_stem}")
            end if
            repeat with d in targetDocs
                save d
            end repeat
            return count of targetDocs
        end tell
        '''
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True,
            timeout=5,
        )
        if result.returncode == 0:
            saved_count = result.stdout.decode().strip()
            if saved_count != "0":
                return True
            else:
                print("⚠️  未找到匹配的演示文稿进行保存")
                return False
        else:
            print(f"⚠️  AppleScript 保存错误: {result.stderr.decode()}")
            return False
    except Exception as e:
        print(f"⚠️  保存演示文稿失败: {e}")
        return False


def close_ppt_document(path: Path) -> bool:
    """
    关闭 PowerPoint 文档（通过 AppleScript）

    Args:
        path: 文档路径

    Returns:
        是否成功
    """
    if platform.system() != "Darwin":
        return False

    try:
        doc_name = path.name
        doc_stem = path.stem
        script = f'''
        tell application "Microsoft PowerPoint"
            set targetDocs to (every presentation whose name is "{doc_name}")
            if (count of targetDocs) is 0 then
                set targetDocs to (every presentation whose name starts with "{doc_stem}")
            end if
            repeat with d in targetDocs
                close d saving no
            end repeat
            return count of targetDocs
        end tell
        '''
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True,
            timeout=5,
        )
        if result.returncode == 0:
            closed_count = result.stdout.decode().strip()
            print(f"📕 已关闭演示文稿: {doc_name} ({closed_count} matched)")
            return True
        else:
            print(f"⚠️  AppleScript 返回错误: {result.stderr.decode()}")
            return False
    except Exception as e:
        print(f"⚠️  关闭演示文稿失败: {e}")
        return False


async def activate_ppt_addin(addin_name: str = DEFAULT_PPT_ADDIN_NAME) -> bool:
    """
    通过 AppleScript UI 自动化激活 PowerPoint Add-In（仅 macOS）

    Args:
        addin_name: Add-In 名称

    Returns:
        是否成功激活
    """
    if platform.system() != "Darwin":
        print("⚠️  自动激活 Add-In 仅支持 macOS，请手动激活")
        return False

    print(f"🔌 尝试自动激活 PPT Add-In: {addin_name}...")

    script_phase1 = '''
tell application "Microsoft PowerPoint" to activate
delay 0.5

tell application "System Events"
    tell process "Microsoft PowerPoint"
        set tg to tab group 1 of front window

        -- 检查功能区是否展开（scroll area 存在则展开）
        set ribbonExpanded to false
        try
            set saCount to count of scroll areas of tg
            if saCount > 0 then set ribbonExpanded to true
        end try

        -- 如果功能区折叠，点击「开始」展开
        if not ribbonExpanded then
            repeat with rb in radio buttons of tg
                if name of rb is "开始" or name of rb is "Home" then
                    click rb
                    delay 0.5
                    exit repeat
                end if
            end repeat
        end if

        -- 使用 entire contents 扁平搜索「加载项」按钮
        set allElems to entire contents of tg
        repeat with elem in allElems
            try
                if role of elem is "AXButton" then
                    set eName to name of elem
                    if eName contains "加载项" or eName contains "Add-in" then
                        click elem
                        return "ok"
                    end if
                end if
            end try
        end repeat

        return "not_found"
    end tell
end tell
'''

    ok, output = _run_applescript(script_phase1)
    if not ok or output != "ok":
        print(f"   ⚠️  未能自动点击「加载项」按钮 (原因: {output})")
        print(f"   👆 请手动点击「加载项」→「{addin_name}...」")
        return False

    print("   ✅ Phase 1: 已点击「加载项」按钮")
    print(f"   👆 请在弹出的加载项面板中点击「{addin_name}...」")
    return False


def dump_ppt_ui_hierarchy() -> None:
    """
    调试工具：打印 PowerPoint 前端窗口的 UI 元素层级

    Usage:
        python -c "from manual_tests.ppt.e2e_base import dump_ppt_ui_hierarchy; dump_ppt_ui_hierarchy()"
    """
    if platform.system() != "Darwin":
        print("仅支持 macOS")
        return

    script = '''
tell application "Microsoft PowerPoint" to activate
delay 0.5

tell application "System Events"
    tell process "Microsoft PowerPoint"
        set output to ""
        set uiElems to UI elements of front window
        repeat with elem in uiElems
            try
                set elemRole to role of elem
                set elemName to name of elem
                set elemDesc to description of elem
                set output to output & "[" & elemRole & "] name=" & elemName & " desc=" & elemDesc & linefeed
                try
                    set children to UI elements of elem
                    repeat with child in children
                        try
                            set cRole to role of child
                            set cName to name of child
                            set cDesc to description of child
                            set output to output & "  [" & cRole & "] name=" & cName & " desc=" & cDesc & linefeed
                        end try
                    end repeat
                end try
            end try
        end repeat
        return output
    end tell
end tell
'''

    ok, output = _run_applescript(script, timeout=30.0)
    if ok:
        print("=== PowerPoint UI Hierarchy ===")
        print(output)
    else:
        print(f"获取 UI 层级失败: {output}")


# ==============================================================================
# PPTTestRunner
# ==============================================================================


class PPTTestRunner(E2ETestRunner):
    """
    PPT E2E 测试运行器

    继承 E2ETestRunner，覆盖 Add-In 激活和文档关闭逻辑以适配 PowerPoint。
    """

    def __init__(
        self,
        fixtures_dir: str | Path | None = None,
        host: str = "127.0.0.1",
        port: int = 3000,
        connection_timeout: float = 30.0,
        auto_open: bool = True,
        auto_close: bool = True,
        auto_activate: bool = True,
        cleanup_on_success: bool = True,
    ):
        super().__init__(
            fixtures_dir=fixtures_dir,
            host=host,
            port=port,
            connection_timeout=connection_timeout,
            auto_open=auto_open,
            auto_close=auto_close,
            auto_activate=auto_activate,
            cleanup_on_success=cleanup_on_success,
        )

    @asynccontextmanager
    async def run_with_workspace(
        self,
        fixture_name: str,
        open_delay: float = 2.0,
    ) -> AsyncIterator[tuple[OfficeWorkspace, DocumentFixture]]:
        """
        准备 PPT 文档并启动 Workspace

        覆盖父类方法以使用 PPT 专属的 Add-In 激活和文档关闭逻辑。
        """
        async with self.prepare_document(fixture_name, open_delay) as fixture:  # pyright: ignore[reportGeneralTypeIssues]
            workspace = OfficeWorkspace(host=self.host, port=self.port)
            try:
                await workspace.start()
                print("✅ Workspace 启动成功")

                # 自动激活 PPT Add-In
                if self.auto_open and self.auto_activate:
                    activated = await activate_ppt_addin()
                    if not activated:
                        print(f"👆 请手动点击「加载项」→「{DEFAULT_PPT_ADDIN_NAME}...」激活 Add-In")

                # 等待 Add-In 连接
                print("⏳ 等待 PPT Add-In 连接...")
                connected = await workspace.wait_for_addin_connection(timeout=self.connection_timeout)
                if not connected:
                    raise RuntimeError("超时：未检测到 PPT Add-In 连接")
                print("✅ PPT Add-In 已连接")

                yield workspace, fixture

            finally:
                if self.auto_close and not fixture.document_closed and fixture.cleanup_on_success:
                    close_ppt_document(fixture.working_path)
                    fixture.document_closed = True
                    await asyncio.sleep(0.5)
                await workspace.stop()

    def verify_slide_info(
        self,
        actual: dict[str, Any],
        expected: ExpectedSlideInfo,
    ) -> tuple[bool, list[str]]:
        """
        验证幻灯片信息

        Args:
            actual: 实际返回的数据
            expected: 预期数据

        Returns:
            (success, messages): 是否通过和验证消息列表
        """
        messages: list[str] = []
        success = True

        if expected.slide_count is not None:
            actual_sc = actual.get("slideCount", 0)
            diff = abs(actual_sc - expected.slide_count)
            if diff <= expected.slide_count_tolerance:
                messages.append(f"✅ 幻灯片数: {actual_sc} (预期 {expected.slide_count})")
            else:
                messages.append(
                    f"❌ 幻灯片数不匹配: {actual_sc} (预期 {expected.slide_count}, "
                    f"误差 {diff} > 允许 {expected.slide_count_tolerance})"
                )
                success = False

        if expected.current_index is not None:
            actual_ci = actual.get("currentSlideIndex", -1)
            if actual_ci == expected.current_index:
                messages.append(f"✅ 当前幻灯片索引: {actual_ci}")
            else:
                messages.append(f"❌ 当前幻灯片索引不匹配: {actual_ci} (预期 {expected.current_index})")
                success = False

        if expected.width is not None:
            actual_w = actual.get("width", 0.0)
            diff = abs(actual_w - expected.width)
            if diff <= expected.width_tolerance:
                messages.append(f"✅ 宽度: {actual_w} (预期 {expected.width})")
            else:
                messages.append(f"❌ 宽度不匹配: {actual_w} (预期 {expected.width})")
                success = False

        if expected.height is not None:
            actual_h = actual.get("height", 0.0)
            diff = abs(actual_h - expected.height)
            if diff <= expected.height_tolerance:
                messages.append(f"✅ 高度: {actual_h} (预期 {expected.height})")
            else:
                messages.append(f"❌ 高度不匹配: {actual_h} (预期 {expected.height})")
                success = False

        return success, messages


# ==============================================================================
# 测试夹具工具
# ==============================================================================


def ensure_ppt_fixtures(fixture_dir: Path) -> dict[str, Path]:
    """
    确保 PPT 测试夹具目录存在

    Args:
        fixture_dir: 夹具目录

    Returns:
        夹具文件路径字典
    """
    fixture_dir.mkdir(parents=True, exist_ok=True)

    paths: dict[str, Path] = {}
    for pptx_file in fixture_dir.glob("*.pptx"):
        paths[pptx_file.name] = pptx_file

    return paths
