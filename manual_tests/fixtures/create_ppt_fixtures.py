"""
PPT E2E 测试 Fixture 生成脚本

使用 python-pptx 程序化生成所有 PPT E2E 测试所需的 fixture .pptx 文件。

Usage:
    uv run python manual_tests/fixtures/create_ppt_fixtures.py
    uv run python manual_tests/fixtures/create_ppt_fixtures.py --clean
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

FIXTURES_ROOT = Path(__file__).parent


# ==============================================================================
# PPT Fixture 生成函数
# ==============================================================================


def create_ppt_empty() -> None:
    """创建空白 PPT (1 张空白幻灯片)"""
    fixture_dir = FIXTURES_ROOT / "ppt_e2e"
    fixture_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    # 添加一张空白幻灯片（使用空白布局）
    blank_layout = prs.slide_layouts[6]  # 通常 index 6 是空白布局
    prs.slides.add_slide(blank_layout)
    prs.save(str(fixture_dir / "empty.pptx"))

    print("  ✅ ppt_e2e/empty.pptx")


def create_ppt_simple() -> None:
    """创建简单 PPT (3 张幻灯片，每页有标题+正文)"""
    fixture_dir = FIXTURES_ROOT / "ppt_e2e"
    fixture_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    slides_content = [
        ("第一页标题", "这是第一页的正文内容。用于测试基本的内容获取功能。"),
        ("第二页标题", "第二页包含一些不同的文本。测试多幻灯片场景。"),
        ("第三页标题", "最后一页内容。PPT E2E 测试结束。"),
    ]

    for title_text, body_text in slides_content:
        # 使用标题+正文布局（通常 index 1）
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.placeholders[0]
        title.text = title_text

        body = slide.placeholders[1]
        body.text = body_text

    prs.save(str(fixture_dir / "simple.pptx"))

    print("  ✅ ppt_e2e/simple.pptx (3 slides)")


def create_ppt_multi_slide() -> None:
    """创建多幻灯片 PPT (5 张幻灯片，编号标题)"""
    fixture_dir = FIXTURES_ROOT / "ppt_e2e"
    fixture_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    for i in range(5):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.placeholders[0]
        title.text = f"幻灯片 {i + 1}"

        body = slide.placeholders[1]
        body.text = f"这是第 {i + 1} 张幻灯片的内容。用于幻灯片管理测试。"

    prs.save(str(fixture_dir / "multi_slide.pptx"))

    print("  ✅ ppt_e2e/multi_slide.pptx (5 slides)")


def create_ppt_multi_element() -> None:
    """创建多元素 PPT (1 张幻灯片，含文本框+表格+形状)"""
    fixture_dir = FIXTURES_ROOT / "ppt_e2e"
    fixture_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 文本框
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    tf.text = "这是一个文本框元素"

    # 表格 (3x3)
    table_shape = slide.shapes.add_table(3, 3, Inches(1), Inches(2.5), Inches(5), Inches(1.5))
    table = table_shape.table
    for row_idx in range(3):
        for col_idx in range(3):
            table.cell(row_idx, col_idx).text = f"R{row_idx + 1}C{col_idx + 1}"

    # 矩形形状
    from pptx.enum.shapes import MSO_SHAPE

    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1), Inches(2), Inches(1.5))

    prs.save(str(fixture_dir / "multi_element.pptx"))

    print("  ✅ ppt_e2e/multi_element.pptx (1 slide, 3 elements)")


def create_ppt_colored() -> None:
    """创建彩色形状 PPT (1 张幻灯片，多个彩色形状，用于截图测试)"""
    fixture_dir = FIXTURES_ROOT / "ppt_e2e"
    fixture_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    shapes_config = [
        (MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1.5), RGBColor(0xFF, 0x00, 0x00)),
        (MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(2), Inches(1.5), RGBColor(0x00, 0xFF, 0x00)),
        (MSO_SHAPE.DIAMOND, Inches(7), Inches(1), Inches(2), Inches(1.5), RGBColor(0x00, 0x00, 0xFF)),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.5), Inches(2), Inches(1.5), RGBColor(0xFF, 0xFF, 0x00)),
        (MSO_SHAPE.PENTAGON, Inches(5.5), Inches(3.5), Inches(2), Inches(1.5), RGBColor(0xFF, 0x00, 0xFF)),
    ]

    for shape_type, left, top, width, height, color in shapes_config:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    prs.save(str(fixture_dir / "colored_slide.pptx"))

    print("  ✅ ppt_e2e/colored_slide.pptx (1 slide, 5 colored shapes)")


# ==============================================================================
# Main
# ==============================================================================


def create_all_ppt_fixtures() -> None:
    """生成所有 PPT fixture 文件"""
    create_ppt_empty()
    create_ppt_simple()
    create_ppt_multi_slide()
    create_ppt_multi_element()
    create_ppt_colored()


def main() -> None:
    """命令行入口"""
    import argparse

    parser = argparse.ArgumentParser(description="生成 PPT E2E 测试 fixture 文件")
    parser.add_argument("--clean", action="store_true", help="清理后重新生成")
    args = parser.parse_args()

    fixture_dir = FIXTURES_ROOT / "ppt_e2e"

    if args.clean and fixture_dir.exists():
        for pptx_file in fixture_dir.glob("*.pptx"):
            pptx_file.unlink()
            print(f"  🗑️  删除: ppt_e2e/{pptx_file.name}")

    print("📝 生成 PPT E2E 测试 fixture 文件...\n")

    create_all_ppt_fixtures()

    print("\n✅ PPT fixture 全部完成")


if __name__ == "__main__":
    main()
